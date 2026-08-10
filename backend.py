import os
import mimetypes
import io
import cv2
import numpy as np
import zipfile
import PyPDF2
import pymupdf as fitz  # PyMuPDF
import uuid
import subprocess
from pptx import Presentation
import json
from typing import List, Dict, Any
from google.genai import types
from fastapi import FastAPI, UploadFile, File, Form, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import StreamingResponse, FileResponse
from pydantic import BaseModel
from google import genai

# --- NEW IMPORTS FOR WORD & EXCEL ---
from pdf2docx import Converter
import pdfplumber
import pandas as pd

# --- Configuration ---
GEMINI_KEY = os.environ.get("GEMINI_API_KEY", "AQ.Ab8RN6IZOita0-V6iAgoHDIVqgLYrfOXdprS1EkMxtml5po2VA")

if os.environ.get("GEMINI_API_KEY") is None:
    client = genai.Client(api_key=GEMINI_KEY)
else:
    client = genai.Client()

app = FastAPI()

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"], 
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# ==========================================
# SYSTEM PROMPT FOR AI APP CONTROL & MEMORY
# ==========================================
APP_CONTROL_PROMPT = """You are an AI assistant built directly into a mobile app called 'PDF Toolkit'. You have the ability to navigate the app and control settings for the user.

IMPORTANT MEMORY RULE: You HAVE access to the user's previous chat history in this conversation. You CAN remember previous questions and context. NEVER say you do not have the ability to recall past conversations. Always use the provided conversation history to answer follow-up questions.

You MUST ALWAYS respond with a valid JSON object. Do not wrap it in markdown block quotes (no ```json).

If the user asks a normal question or needs a PDF summary, respond with JSON like this:
{"reply": "Your normal conversational response here."}

HOWEVER, if the user asks you to navigate the app, change a setting, or open a tool, you must also include the 'action' and 'target' fields in the JSON. 

Exact commands you can use:
- To go to the Profile screen: {"reply": "Opening your profile now.", "action": "NAVIGATE", "target": "3"}
- To go to the Search screen: {"reply": "Opening the search tab.", "action": "NAVIGATE", "target": "1"}
- To go to the Home/Dashboard screen: {"reply": "Going to the dashboard.", "action": "NAVIGATE", "target": "0"}
- To turn on/off Dark Mode: {"reply": "Toggling dark mode for you.", "action": "TOGGLE_DARK_MODE"}
- To clear the chat history: {"reply": "Wiping chat history.", "action": "CLEAR_CHAT"}
- To open a specific tool (e.g., 'Merge PDFs', 'PDF to JPG', 'Add Watermark'): {"reply": "Opening tool.", "action": "OPEN_TOOL", "target": "Exact Tool Name"}
"""

@app.get("/")
@app.head("/")
def read_root():
    return {"message": "PDF Toolkit Backend is running!"}

# ==========================================
# ENDPOINT: AI PDF SUMMARY (Single File)
# ==========================================
@app.post("/api/summarize")
async def summarize_pdf(file: UploadFile = File(...)):
    try:
        text = ""
        pdf_reader = PyPDF2.PdfReader(file.file)
        for page in pdf_reader.pages:
            extracted = page.extract_text()
            if extracted:
                text += extracted + "\n"
        
        if not text.strip():
            return {"status": "error", "message": "Could not extract any text from this PDF."}

        prompt = f"Provide a clear, well-structured summary of this document using bullet points:\n\n{text}"
        response = client.models.generate_content(model='gemini-2.5-flash', contents=prompt)
        
        return {"status": "success", "result": response.text}
        
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# ENDPOINT: ADD PASSWORD TO PDF (Single File)
# ==========================================
@app.post("/api/lock")
async def lock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    try:
        reader = PyPDF2.PdfReader(file.file)
        writer = PyPDF2.PdfWriter()
        
        for page in reader.pages:
            writer.add_page(page)
            
        writer.encrypt(password)
        
        memory_file = io.BytesIO()
        writer.write(memory_file)
        memory_file.seek(0)
            
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="locked_{file.filename}"'}
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# ENDPOINT: PDF TO JPG (Single File -> ZIP output)
# ==========================================
@app.post("/api/tojpg")
async def pdf_to_jpg(file: UploadFile = File(...)):
    try:
        file_bytes = await file.read()
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        
        if len(doc) == 0:
            return {"status": "error", "message": "The uploaded PDF is completely empty."}
            
        memory_zip = io.BytesIO()
        with zipfile.ZipFile(memory_zip, "w") as zf:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                pix = page.get_pixmap(dpi=150)
                img_bytes = pix.tobytes("jpg")
                zf.writestr(f"page_{page_num + 1}.jpg", img_bytes)
        
        memory_zip.seek(0)
        
        return StreamingResponse(
            memory_zip, 
            media_type='application/zip',
            headers={'Content-Disposition': f'attachment; filename="converted_{file.filename}.zip"'}
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: MERGE PDFs (Multiple Files)
# ==========================================
@app.post("/api/merge")
async def merge_pdfs(files: List[UploadFile] = File(...)):
    """Receives multiple PDFs and combines them into one."""
    try:
        print(f"📑 Merging {len(files)} PDFs...")
        merger = PyPDF2.PdfMerger()
        
        for file in files:
            file_bytes = await file.read()
            merger.append(io.BytesIO(file_bytes))
            
        memory_file = io.BytesIO()
        merger.write(memory_file)
        memory_file.seek(0)
        
        print("✅ Merge successful! Sending back to app...")
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': 'attachment; filename="merged_document.pdf"'}
        )
    except Exception as e:
        print(f"⚠️ Error: {e}")
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: JPG TO PDF (Multiple Files)
# ==========================================
@app.post("/api/topdf")
async def jpg_to_pdf(files: List[UploadFile] = File(...)):
    """Receives multiple images and packs them into a single PDF."""
    try:
        print(f"🖼️ Packing {len(files)} images into PDF...")
        doc = fitz.open() # Create a blank PDF document
        
        for file in files:
            img_bytes = await file.read()
            
            # Determine extension based on file name (jpg, png, etc.)
            ext = file.filename.split('.')[-1].lower() if '.' in file.filename else 'jpg'
            
            # Open the image with PyMuPDF and convert directly to PDF bytes
            img_doc = fitz.open(stream=img_bytes, filetype=ext)
            pdf_bytes = img_doc.convert_to_pdf()
            img_doc.close()
            
            # Insert the newly converted page into the main document
            temp_pdf = fitz.open("pdf", pdf_bytes)
            doc.insert_pdf(temp_pdf)
            temp_pdf.close()
            
        memory_file = io.BytesIO()
        doc.save(memory_file)
        memory_file.seek(0)
        
        print("✅ Conversion successful! Sending back to app...")
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': 'attachment; filename="converted_images.pdf"'}
        )
    except Exception as e:
        print(f"⚠️ Error: {e}")
        return {"status": "error", "message": str(e)}
    
# ==========================================
# ENDPOINT: REMOVE PASSWORD FROM PDF
# ==========================================
@app.post("/api/unlock")
async def unlock_pdf(file: UploadFile = File(...), password: str = Form(...)):
    """Receives an encrypted PDF and a password, decrypts it, and returns the unlocked file."""
    try:
        print(f"🔓 Decrypting file: {file.filename}")
        
        reader = PyPDF2.PdfReader(file.file)
        
        # 1. Check if the file is actually encrypted
        if not reader.is_encrypted:
            return {"status": "error", "message": "This PDF is already unlocked!"}
            
        # 2. Attempt to decrypt with the provided password
        success = reader.decrypt(password)
        
        if success == 0:
            return {"status": "error", "message": "Incorrect password. Please try again."}
            
        # 3. If successful, copy the decrypted pages to a new writer
        writer = PyPDF2.PdfWriter()
        for page in reader.pages:
            writer.add_page(page)
            
        # 4. Save to memory and send back
        memory_file = io.BytesIO()
        writer.write(memory_file)
        memory_file.seek(0)
        
        print("✅ PDF Decrypted successfully! Sending back to app...")
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="unlocked_{file.filename}"'}
        )
        
    except Exception as e:
        print(f"⚠️ Error: {e}")
        return {"status": "error", "message": str(e)}

class ChatRequest(BaseModel):
    message: str
    history: List[Dict[str, str]] = [] # Accepts the chat history

# ==========================================
# ENDPOINT: AI CHAT (APP CONTROL & MEMORY INTEGRATED)
# ==========================================
@app.post("/api/chat")
async def chat_with_ai(request: ChatRequest):
    try:
        my_api_key = os.environ.get("GEMINI_API_KEY")
        client = genai.Client(api_key=my_api_key)
        
        # 1. Build the conversation history for Gemini
        contents = []
        for msg in request.history:
            role = "user" if msg.get("sender") == "user" else "model"
            contents.append(
                types.Content(role=role, parts=[types.Part.from_text(text=msg.get("text", ""))])
            )
            
        # 2. Add the current message
        contents.append(
            types.Content(role="user", parts=[types.Part.from_text(text=request.message)])
        )
        
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=contents,
            config=types.GenerateContentConfig(
                system_instruction=APP_CONTROL_PROMPT,
                response_mime_type="application/json",
            )
        )
        
        ai_data = json.loads(response.text)
        
        return {
            "status": "success", 
            "reply": ai_data.get("reply", "I am not sure how to respond to that."),
            "action": ai_data.get("action"),
            "target": ai_data.get("target")
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# ENDPOINT: AI CHAT (FILE ATTACHED WITH MEMORY)
# ==========================================
@app.post("/api/chat/file")
async def chat_with_file(
    file: UploadFile = File(...), 
    message: str = Form("Please summarize this document in detail."),
    history: str = Form("[]") # Receives history as a JSON string
):
    try:
        my_api_key = os.environ.get("GEMINI_API_KEY")
        client = genai.Client(api_key=my_api_key)
        
        # Extract text from the uploaded PDF
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        extracted_text = ""
        for page in doc:
            extracted_text += page.get_text()
            
        # Build the conversation history for Gemini
        contents = []
        parsed_history = json.loads(history)
        for msg in parsed_history:
            role = "user" if msg.get("sender") == "user" else "model"
            contents.append(
                types.Content(role=role, parts=[types.Part.from_text(text=msg.get("text", ""))])
            )
            
        # Combine the user's prompt with the document text
        full_prompt = f"{message}\n\n--- Document Content ---\n{extracted_text}"
        contents.append(
            types.Content(role="user", parts=[types.Part.from_text(text=full_prompt)])
        )
        
        # Send to Gemini
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=contents,
            config=types.GenerateContentConfig(
                system_instruction=APP_CONTROL_PROMPT,
                response_mime_type="application/json",
            )
        )
        
        ai_data = json.loads(response.text)
        
        return {
            "status": "success", 
            "reply": ai_data.get("reply", "I have processed the document."),
            "action": ai_data.get("action"),
            "target": ai_data.get("target")
        }
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# 1. PDF METADATA INSPECTOR
# ==========================================
@app.post("/api/metadata")
async def extract_metadata(file: UploadFile = File(...)):
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        metadata = doc.metadata
        
        result_text = "--- PDF Metadata ---\n"
        for key, value in metadata.items():
            if value: 
                result_text += f"{key.capitalize()}: {value}\n"
        
        if result_text == "--- PDF Metadata ---\n":
            result_text = "No metadata found in this document."
            
        return {"status": "success", "result": result_text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# 2. DOCUMENT WATERMARKING
# ==========================================
@app.post("/api/watermark")
async def add_watermark(
    file: UploadFile = File(...), 
    password: str = Form("CONFIDENTIAL") 
):
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        for page in doc:
            rect = page.rect
            point = fitz.Point(rect.width * 0.15, rect.height * 0.50)
            page.insert_text(
                point, 
                password, 
                fontsize=60, 
                color=(0.7, 0.7, 0.7) 
            )
            
        out_bytes = doc.write()
        return Response(content=out_bytes, media_type="application/pdf")
    except Exception as e:
        return {"status": "error", "message": str(e)}
        
# ==========================================
# 3. AI OPTICAL CHARACTER RECOGNITION (OCR)
# ==========================================
@app.post("/api/imagetotext")
async def image_to_text(file: UploadFile = File(...)):
    try:
        image_bytes = await file.read()
        my_api_key = os.environ.get("GEMINI_API_KEY")
        client = genai.Client(api_key=my_api_key)
        
        mime_type, _ = mimetypes.guess_type(file.filename)
        
        if not mime_type or mime_type == "application/octet-stream":
            if file.filename.lower().endswith('.png'):
                mime_type = "image/png"
            else:
                mime_type = "image/jpeg"
        
        image_part = types.Part.from_bytes(
            data=image_bytes,
            mime_type=mime_type,
        )
        
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=[
                image_part,
                "Extract all the text from this image exactly as it appears. Do not add any extra commentary."
            ]
        )
        return {"status": "success", "result": response.text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: COMPRESS PDF
# ==========================================
@app.post("/api/compress")
async def compress_pdf(file: UploadFile = File(...)):
    """Receives a PDF, compresses its streams and removes unused data, and returns it."""
    try:
        print(f"🗜️ Compressing file: {file.filename}")
        
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        memory_file = io.BytesIO()
        doc.save(memory_file, garbage=4, deflate=True, deflate_images=True)
        memory_file.seek(0)
        
        print("✅ Compression successful! Sending back to app...")
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="compressed_{file.filename}"'}
        )
        
    except Exception as e:
        print(f"⚠️ Error: {e}")
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: ROTATE PDF
# ==========================================
@app.post("/api/rotate")
async def rotate_pdf(file: UploadFile = File(...), input_data: str = Form("90")):
    try:
        degrees = int(input_data.strip())
        reader = PyPDF2.PdfReader(file.file)
        writer = PyPDF2.PdfWriter()
        
        for page in reader.pages:
            page.rotate(degrees)
            writer.add_page(page)
            
        memory_file = io.BytesIO()
        writer.write(memory_file)
        memory_file.seek(0)
        
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="rotated_{file.filename}"'}
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: SPLIT PDF
# ==========================================
@app.post("/api/split")
async def split_pdf(file: UploadFile = File(...), input_data: str = Form("1-1")):
    try:
        parts = input_data.split('-')
        start_page = max(1, int(parts[0].strip()))
        end_page = int(parts[1].strip()) if len(parts) > 1 else start_page
        
        reader = PyPDF2.PdfReader(file.file)
        writer = PyPDF2.PdfWriter()
        
        start_idx = start_page - 1
        end_idx = min(len(reader.pages), end_page)
        
        for i in range(start_idx, end_idx):
            writer.add_page(reader.pages[i])
            
        memory_file = io.BytesIO()
        writer.write(memory_file)
        memory_file.seek(0)
        
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="split_{file.filename}"'}
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: EXTRACT IMAGES FROM PDF
# ==========================================
@app.post("/api/extractimages")
async def extract_images(file: UploadFile = File(...)):
    try:
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        
        memory_zip = io.BytesIO()
        image_count = 0
        
        with zipfile.ZipFile(memory_zip, "w") as zf:
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                image_list = page.get_images(full=True)
                
                for img_index, img in enumerate(image_list):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    ext = base_image["ext"]
                    image_count += 1
                    zf.writestr(f"image_p{page_num+1}_{image_count}.{ext}", image_bytes)
        
        if image_count == 0:
            return {"status": "error", "message": "No images found in this PDF."}
            
        memory_zip.seek(0)
        return StreamingResponse(
            memory_zip, 
            media_type='application/zip',
            headers={'Content-Disposition': f'attachment; filename="extracted_images_{file.filename}.zip"'}
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: AI PDF TRANSLATOR
# ==========================================
@app.post("/api/translate")
async def translate_pdf(file: UploadFile = File(...), input_data: str = Form("Hindi")):
    try:
        my_api_key = os.environ.get("GEMINI_API_KEY")
        client = genai.Client(api_key=my_api_key)
        
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        extracted_text = ""
        for page in doc:
            extracted_text += page.get_text()
            
        if not extracted_text.strip():
             return {"status": "error", "message": "No text found to translate."}

        prompt = f"Translate the following document into {input_data}. Maintain the original tone and structure as much as possible:\n\n{extracted_text}"
        
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt
        )
        
        return {"status": "success", "result": response.text}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: PDF TO WORD
# ==========================================
@app.post("/api/toword")
async def pdf_to_word(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    try:
        # Create unique temporary filenames
        pdf_path = f"temp_{uuid.uuid4().hex}_{file.filename}"
        docx_path = pdf_path.replace('.pdf', '.docx')
        
        # Save uploaded PDF to disk
        with open(pdf_path, "wb") as buffer:
            buffer.write(await file.read())
        
        # Convert using pdf2docx
        cv = Converter(pdf_path)
        cv.convert(docx_path)
        cv.close()
        
        # Define cleanup function to delete temp files after sending
        def cleanup():
            if os.path.exists(pdf_path): os.remove(pdf_path)
            if os.path.exists(docx_path): os.remove(docx_path)
            
        background_tasks.add_task(cleanup)
        
        return FileResponse(
            path=docx_path, 
            filename=f"converted_{file.filename.replace('.pdf', '.docx')}",
            media_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINT: PDF TO EXCEL
# ==========================================
@app.post("/api/toexcel")
async def pdf_to_excel(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    try:
        # Create unique temporary filenames
        pdf_path = f"temp_{uuid.uuid4().hex}_{file.filename}"
        excel_path = pdf_path.replace('.pdf', '.xlsx')
        
        # Save uploaded PDF to disk
        with open(pdf_path, "wb") as buffer:
            buffer.write(await file.read())
            
        all_tables = []
        
        # Extract tables using pdfplumber
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                tables = page.extract_tables()
                for table in tables:
                    if table:
                        df = pd.DataFrame(table[1:], columns=table[0])
                        all_tables.append(df)
        
        if not all_tables:
            if os.path.exists(pdf_path): os.remove(pdf_path)
            return {"status": "error", "message": "No tabular data found in this PDF"}
            
        # Combine all extracted tables into one dataframe and save as excel
        final_df = pd.concat(all_tables, ignore_index=True)
        final_df.to_excel(excel_path, index=False)
        
        # Define cleanup function to delete temp files after sending
        def cleanup():
            if os.path.exists(pdf_path): os.remove(pdf_path)
            if os.path.exists(excel_path): os.remove(excel_path)
            
        background_tasks.add_task(cleanup)
        
        return FileResponse(
            path=excel_path, 
            filename=f"converted_{file.filename.replace('.pdf', '.xlsx')}",
            media_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet'
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host='0.0.0.0', port=5000)

# ==========================================
# NEW ENDPOINT: PDF TO POWERPOINT (PPTX)
# ==========================================
@app.post("/api/pdftoppt")
async def pdf_to_ppt(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    """Converts a PDF into a PowerPoint by turning each page into an image slide."""
    try:
        pdf_path = f"temp_{uuid.uuid4().hex}_{file.filename}"
        pptx_path = pdf_path.replace('.pdf', '.pptx')
        
        with open(pdf_path, "wb") as buffer:
            buffer.write(await file.read())
            
        doc = fitz.open(pdf_path)
        prs = Presentation()
        blank_slide_layout = prs.slide_layouts[6] # 6 is the layout for a blank slide
        
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150)
            img_path = f"temp_img_{uuid.uuid4().hex}.png"
            pix.save(img_path)
            
            slide = prs.slides.add_slide(blank_slide_layout)
            # Add image to fill the entire slide
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)
            os.remove(img_path)
            
        prs.save(pptx_path)
        
        def cleanup():
            if os.path.exists(pdf_path): os.remove(pdf_path)
            if os.path.exists(pptx_path): os.remove(pptx_path)
            
        background_tasks.add_task(cleanup)
        
        return FileResponse(
            path=pptx_path, 
            filename=f"converted_{file.filename.replace('.pdf', '.pptx')}",
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation'
        )
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# NEW ENDPOINTS: OFFICE (Word/Excel/PPT) TO PDF
# ==========================================
async def convert_office_to_pdf_logic(background_tasks, file, extension):
    """Core logic using LibreOffice headless to convert any office format to PDF."""
    try:
        input_path = f"temp_{uuid.uuid4().hex}_{file.filename}"
        output_dir = os.path.dirname(os.path.abspath(input_path)) or "."
        
        with open(input_path, "wb") as buffer:
            buffer.write(await file.read())
            
        # Execute LibreOffice headless conversion
        subprocess.run(["libreoffice", "--headless", "--convert-to", "pdf", input_path, "--outdir", output_dir], check=True)
        
        pdf_path = input_path.rsplit('.', 1)[0] + ".pdf"
        
        def cleanup():
            if os.path.exists(input_path): os.remove(input_path)
            if os.path.exists(pdf_path): os.remove(pdf_path)
            
        background_tasks.add_task(cleanup)
        
        return FileResponse(
            path=pdf_path, 
            filename=f"converted_{file.filename.rsplit('.', 1)[0]}.pdf",
            media_type='application/pdf'
        )
    except Exception as e:
        return {"status": "error", "message": f"Server requires LibreOffice. Error: {str(e)}"}

@app.post("/api/wordtopdf")
async def word_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    return await convert_office_to_pdf_logic(background_tasks, file, "word")

@app.post("/api/exceltopdf")
async def excel_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    return await convert_office_to_pdf_logic(background_tasks, file, "excel")

@app.post("/api/ppttopdf")
async def ppt_to_pdf(background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    return await convert_office_to_pdf_logic(background_tasks, file, "ppt")

# ==========================================
# ENDPOINT: AI QUIZ GENERATOR
# ==========================================
@app.post("/api/quiz")
async def generate_quiz(file: UploadFile = File(...), num_questions: str = Form("5")):
    try:
        import os
        import fitz
        import json
        from google import genai
        
        my_api_key = os.environ.get("GEMINI_API_KEY")
        client = genai.Client(api_key=my_api_key)
        
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        extracted_text = ""
        for page in doc:
            extracted_text += page.get_text()
            
        if not extracted_text.strip():
             return {"status": "error", "message": "No text found to generate a quiz."}

        # Dynamically inject the requested number of questions into the prompt
        prompt = f"""Generate a {num_questions}-question multiple choice quiz based on the following text.
        Return ONLY a valid JSON array of objects. Do not include any markdown formatting like ```json.
        Strict format required:
        [
          {{
            "question": "Question text?",
            "options": ["Option A", "Option B", "Option C", "Option D"],
            "answer": "Option A",
            "explanation": "Brief explanation of why Option A is correct."
          }}
        ]

        Text to quiz on:
        {extracted_text[:20000]}""" 
        
        response = client.models.generate_content(
            model='gemini-2.5-flash',
            contents=prompt
        )
        
        # Clean up response if Gemini ignores the "no markdown" rule
        res_text = response.text.strip()
        if res_text.startswith("```json"):
            res_text = res_text[7:-3]
        elif res_text.startswith("```"):
            res_text = res_text[3:-3]

        quiz_data = json.loads(res_text.strip())
        
        return {"status": "success", "result": quiz_data}
    except Exception as e:
        return {"status": "error", "message": str(e)}

# ==========================================
# ENDPOINT: AI SCAN ENHANCER (OPENCV)
# ==========================================
@app.post("/api/enhance-scan")
async def enhance_scan(background_tasks: BackgroundTasks, file: UploadFile = File(...), filter_type: str = Form(...)):
    """Receives a scanned PDF, extracts the images, applies advanced AI/OpenCV filters, and returns a new PDF."""
    try:
        import io
        import cv2
        import numpy as np
        import fitz
        from fastapi.responses import StreamingResponse

        print(f"🪄 Applying {filter_type} filter to {file.filename}...")
        pdf_bytes = await file.read()
        doc = fitz.open(stream=pdf_bytes, filetype="pdf")
        out_doc = fitz.open()

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=200) 
            
            img_data = pix.tobytes("png")
            nparr = np.frombuffer(img_data, np.uint8)
            img = cv2.imdecode(nparr, cv2.IMREAD_COLOR)

            # --- NEW COMBINED OPENCV FILTERS ---
            if filter_type == "Auto":
                lab = cv2.cvtColor(img, cv2.COLOR_BGR2LAB)
                l, a, b = cv2.split(lab)
                clahe = cv2.createCLAHE(clipLimit=2.0, tileGridSize=(8,8))
                l = clahe.apply(l)
                img = cv2.cvtColor(cv2.merge((l,a,b)), cv2.COLOR_LAB2BGR)
                
            elif filter_type == "Color":
                hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV).astype(np.float32)
                hsv[:, :, 1] = hsv[:, :, 1] * 1.3 # Boost saturation by 30%
                hsv[:, :, 1] = np.clip(hsv[:, :, 1], 0, 255)
                img = cv2.cvtColor(hsv.astype(np.uint8), cv2.COLOR_HSV2BGR)
                
            elif filter_type == "Grayscale":
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                img = cv2.cvtColor(gray, cv2.COLOR_GRAY2BGR)
                
            elif filter_type == "B&W":
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                bw = cv2.adaptiveThreshold(gray, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY, 21, 15)
                img = cv2.cvtColor(bw, cv2.COLOR_GRAY2BGR)
                
            elif filter_type == "Shadows":
                hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV).astype(np.float32)
                h, s, v = cv2.split(hsv)
                v = v + (255 - v) * 0.4 # Recover shadows by boosting darker areas
                v = np.clip(v, 0, 255)
                img = cv2.cvtColor(cv2.merge((h, s, v)).astype(np.uint8), cv2.COLOR_HSV2BGR)

            elif filter_type == "Enhance":
                lab = cv2.cvtColor(img, cv2.COLOR_BGR2LAB)
                l_channel, a, b = cv2.split(lab)
                clahe = cv2.createCLAHE(clipLimit=2.5, tileGridSize=(8,8))
                cl = clahe.apply(l_channel)
                img = cv2.cvtColor(cv2.merge((cl,a,b)), cv2.COLOR_LAB2BGR)
                
            elif filter_type == "Soft":
                img = cv2.bilateralFilter(img, 9, 75, 75)
                
            elif filter_type == "Brighten":
                hsv = cv2.cvtColor(img, cv2.COLOR_BGR2HSV)
                h, s, v = cv2.split(hsv)
                v = cv2.add(v, 50)
                v = np.clip(v, 0, 255)
                img = cv2.cvtColor(cv2.merge((h, s, v)), cv2.COLOR_HSV2BGR)
                
            elif filter_type == "Remove Shadow":
                rgb_planes = cv2.split(img)
                result_planes = []
                for plane in rgb_planes:
                    dilated_img = cv2.dilate(plane, np.ones((7,7), np.uint8))
                    bg_img = cv2.medianBlur(dilated_img, 21)
                    diff_img = 255 - cv2.absdiff(plane, bg_img)
                    norm_img = cv2.normalize(diff_img, None, alpha=0, beta=255, norm_type=cv2.NORM_MINMAX, dtype=cv2.CV_8UC1)
                    result_planes.append(norm_img)
                img = cv2.merge(result_planes)
                
            elif filter_type == "Erase Handwriting":
                gray = cv2.cvtColor(img, cv2.COLOR_BGR2GRAY)
                _, thresh = cv2.threshold(gray, 120, 255, cv2.THRESH_BINARY)
                kernel = np.ones((2,2), np.uint8)
                img = cv2.morphologyEx(thresh, cv2.MORPH_OPEN, kernel)
                img = cv2.cvtColor(img, cv2.COLOR_GRAY2BGR)
                
            elif filter_type == "Remove Moire":
                img = cv2.fastNlMeansDenoisingColored(img, None, 10, 10, 7, 21)

            # --- CONVERT BACK TO PDF ---
            is_success, buffer = cv2.imencode(".png", img)
            img_doc = fitz.open(stream=buffer.tobytes(), filetype="png")
            pdf_bytes_page = img_doc.convert_to_pdf()
            img_doc.close()
            
            temp_pdf = fitz.open("pdf", pdf_bytes_page)
            out_doc.insert_pdf(temp_pdf)
            temp_pdf.close()

        memory_file = io.BytesIO()
        out_doc.save(memory_file)
        memory_file.seek(0)
        
        doc.close()
        out_doc.close()

        print("✅ Filter applied successfully!")
        return StreamingResponse(
            memory_file, 
            media_type='application/pdf',
            headers={'Content-Disposition': f'attachment; filename="enhanced_{file.filename}"'}
        )
    except Exception as e:
        print(f"⚠️ Error enhancing scan: {str(e)}")
        return {"status": "error", "message": str(e)}
